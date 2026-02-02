"""
Unified PowerPoint Generator

This module combines template-based and programmatic generation approaches.
It primarily uses a template-driven approach to preserve professional design,
with fallback to programmatic generation if template is not available.

Template-Driven Approach:
- Loads pre-designed template.pptx
- Replaces {{MARKER}} placeholders with actual data
- Inserts images (certifications, partners) from web search
- Preserves all formatting, colors, fonts, and layouts

Why Template-Driven?
- Preserves professional design and formatting
- Ensures brand consistency (colors, fonts, layout)
- Faster generation (no need to programmatically create layouts)
- Easy to update design without changing code
- Designers can work directly in PowerPoint
"""

import os
import re
import requests
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image

# ============================================================
# CONFIGURATION
# ============================================================
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "template.pptx")

# Branding Colors
KELP_INDIGO = RGBColor(45, 0, 75)
KELP_PINK = RGBColor(255, 0, 127)
KELP_ORANGE = RGBColor(255, 100, 50)
KELP_WHITE = RGBColor(255, 255, 255)
KELP_GREY_TEXT = RGBColor(80, 80, 80)
KELP_LIGHT_GREY_BG = RGBColor(245, 245, 245)

# Temporary directory for downloaded images
TEMP_IMAGE_DIR = None

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def clean_text(text):
    """Removes markdown and cleans text for presentation."""
    if not text:
        return ""
    if isinstance(text, (int, float)):
        return str(text)
    text = str(text).replace("**", "").replace("*", "").strip()
    return text

def format_list(items, max_items=None, bullet="•"):
    """Formats a list of items into a bulleted string."""
    if not items:
        return ""
    if isinstance(items, str):
        items = [items]
    if max_items:
        items = items[:max_items]
    return "\n".join([f"{bullet} {clean_text(item)}" for item in items])

def get_unique_output_path(output_path):
    """
    Generates a unique output path by adding a number suffix if file exists.
    Examples: file.pptx -> file_1.pptx -> file_2.pptx
    """
    if not os.path.exists(output_path):
        return output_path
    
    directory = os.path.dirname(output_path)
    base_name = os.path.basename(output_path)
    name_without_ext, ext = os.path.splitext(base_name)
    
    counter = 1
    while True:
        new_name = f"{name_without_ext}_{counter}{ext}"
        new_path = os.path.join(directory, new_name)
        if not os.path.exists(new_path):
            print(f"   ℹ️  File exists, using: {new_name}")
            return new_path
        counter += 1
        if counter > 1000:
            raise ValueError("Too many files with similar names exist")

def download_image_to_temp(image_url: str, filename: str) -> str:
    """Downloads an image from URL and saves it to a temporary file."""
    global TEMP_IMAGE_DIR
    
    if not TEMP_IMAGE_DIR:
        TEMP_IMAGE_DIR = tempfile.mkdtemp(prefix="kelp_images_")
    
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(image_url, headers=headers, timeout=10, stream=True)
        response.raise_for_status()
        
        file_path = os.path.join(TEMP_IMAGE_DIR, filename)
        with open(file_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        
        # Validate it's a valid image
        img = Image.open(file_path)
        img.verify()
        return file_path
    except Exception as e:
        print(f"   ⚠️  Failed to download image {image_url}: {e}")
        return None

# ============================================================
# TEMPLATE REPLACEMENT FUNCTIONS (Using Placeholder Indices)
# ============================================================

def fill_placeholder_by_index(slide, placeholder_idx, content, content_type='text'):
    """
    Fills a placeholder by its index with the given content.
    
    Args:
        slide: PowerPoint slide object
        placeholder_idx: Index of the placeholder (e.g., 10, 11, 12)
        content: Content to fill (text string, image path, etc.)
        content_type: Type of content - 'text', 'image', 'chart'
    
    Returns:
        True if successful, False otherwise
    """
    try:
        # Check if placeholder exists
        if placeholder_idx not in slide.placeholders:
            return False
        
        placeholder = slide.placeholders[placeholder_idx]
        placeholder_type = placeholder.placeholder_format.type
        
        if content_type == 'text':
            # Fill text content into BODY or OBJECT placeholders
            if placeholder_type in [2, 7]:  # BODY (2) or OBJECT (7)
                if hasattr(placeholder, 'text_frame'):
                    tf = placeholder.text_frame
                    tf.clear()
                    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                    p.text = content
                    return True
        
        elif content_type == 'image':
            # Fill image into PICTURE placeholder
            if placeholder_type == 18:  # PICTURE (18)
                if isinstance(content, str) and os.path.exists(content):
                    # Insert image into placeholder
                    placeholder.insert_picture(content)
                    return True
        
        return False
    except Exception as e:
        print(f"   ⚠️  Failed to fill placeholder {placeholder_idx}: {e}")
        return False

def detect_placeholder_info(slide, placeholder_idx):
    """
    Detects placeholder format information similar to index.py.
    Returns dict with placeholder details or None if not found.
    """
    try:
        placeholder = slide.placeholders[placeholder_idx]
        return {
            'idx': placeholder.placeholder_format.idx,
            'type': placeholder.placeholder_format.type,
            'name': placeholder.name,
            'placeholder': placeholder
        }
    except (KeyError, AttributeError):
        return None

def calculate_optimal_font_size(text_length, placeholder_width, placeholder_height, min_size=8, max_size=18):
    """
    Calculates optimal font size based on text length and placeholder dimensions.
    Prevents text overflow by reducing font size if needed.
    """
    # Estimate characters per line based on placeholder width
    # Rough estimate: 1 inch ≈ 12 characters at 12pt font
    chars_per_line = int((placeholder_width / Inches(1)) * 12)
    lines_available = int((placeholder_height / Inches(1)) * 2)  # Rough estimate
    
    if chars_per_line == 0 or lines_available == 0:
        return 12  # Default size
    
    # Calculate how many lines the text will need
    estimated_lines = max(1, text_length / chars_per_line)
    
    # If text fits comfortably, use max size
    if estimated_lines <= lines_available * 0.8:
        return max_size
    
    # If text is too long, reduce font size proportionally
    ratio = (lines_available * 0.8) / estimated_lines
    calculated_size = int(max_size * ratio)
    
    # Clamp between min and max
    return max(min_size, min(calculated_size, max_size))

def fill_text_placeholder(slide, placeholder_idx, text_content):
    """
    Helper to fill a text placeholder with formatted text.
    Detects placeholder format and adjusts font size to prevent overflow.
    """
    try:
        # Detect placeholder info (like index.py)
        placeholder_info = detect_placeholder_info(slide, placeholder_idx)
        if not placeholder_info:
            return False
        
        placeholder = placeholder_info['placeholder']
        placeholder_type = placeholder_info['type']
        placeholder_name = placeholder_info['name']
        
        # Only fill BODY (2) or OBJECT (7) placeholders
        if placeholder_type not in [2, 7]:
            print(f"   ⚠️  Placeholder {placeholder_idx} is type {placeholder_type}, not a text placeholder")
            return False
        
        if not hasattr(placeholder, 'text_frame'):
            return False
        
        tf = placeholder.text_frame
        tf.clear()
        tf.word_wrap = True  # Enable word wrapping
        # Set margins to ensure text fits
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        # Get placeholder dimensions for font size calculation
        placeholder_width = placeholder.width
        placeholder_height = placeholder.height
        
        # Calculate optimal font size based on content length
        text_length = len(text_content)
        optimal_font_size = calculate_optimal_font_size(
            text_length, 
            placeholder_width, 
            placeholder_height,
            min_size=8,
            max_size=14  # Reduced max to prevent overflow
        )
        
        # Split text into lines if it's a list format (bullet points)
        if '\n' in text_content:
            lines = text_content.split('\n')
        else:
            # For long text, try to split intelligently
            lines = [text_content]
            if len(text_content) > 200:
                # Split long text into chunks
                words = text_content.split()
                lines = []
                current_line = []
                current_length = 0
                max_chars_per_line = 80
                
                for word in words:
                    if current_length + len(word) + 1 > max_chars_per_line and current_line:
                        lines.append(' '.join(current_line))
                        current_line = [word]
                        current_length = len(word)
                    else:
                        current_line.append(word)
                        current_length += len(word) + 1
                
                if current_line:
                    lines.append(' '.join(current_line))
        
        # Add paragraphs for each line
        for i, line in enumerate(lines):
            if not line.strip():
                continue
                
            if i == 0:
                # Use first paragraph
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            else:
                p = tf.add_paragraph()
            
            # Set text (this creates a run automatically)
            p.text = line.strip()
            
            # Set font properties on the first run
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(optimal_font_size)
                run.font.name = 'Arial'
            
            # Add spacing between paragraphs
            if i < len(lines) - 1:
                p.space_after = Pt(4)
        
        print(f"      📝 Filled placeholder {placeholder_idx} ({placeholder_name}, type {placeholder_type}) with font size {optimal_font_size}pt")
        return True
        
    except (KeyError, AttributeError) as e:
        # Placeholder doesn't exist or doesn't have text_frame
        return False
    except Exception as e:
        print(f"   ⚠️  Failed to fill text placeholder {placeholder_idx}: {e}")
        return False

def fill_image_placeholder(slide, placeholder_idx, image_path):
    """Helper to fill a picture placeholder with an image."""
    try:
        placeholder = slide.placeholders[placeholder_idx]
        placeholder_type = placeholder.placeholder_format.type
        
        if placeholder_type == 18:  # PICTURE (18)
            if os.path.exists(image_path):
                # Get placeholder position and size
                left = placeholder.left
                top = placeholder.top
                width = placeholder.width
                height = placeholder.height
                
                # Add picture at the same position (this will overlay/replace the placeholder)
                slide.shapes.add_picture(image_path, left, top, width, height)
                return True
        return False
    except (KeyError, AttributeError):
        # Placeholder doesn't exist
        return False
    except Exception as e:
        print(f"   ⚠️  Failed to fill image placeholder {placeholder_idx}: {e}")
        return False

# ============================================================
# IMAGE INSERTION FUNCTIONS
# ============================================================

def insert_customer_images_by_placeholder(slide, partners_data, placeholder_indices=[11, 12, 13]):
    """
    Inserts customer/partner images into picture placeholders (indices 11, 12, 13).
    """
    if not partners_data:
        return
    
    # Extract partner information
    partners = []
    if isinstance(partners_data, list):
        partners = partners_data
    elif isinstance(partners_data, dict) and 'partners' in partners_data:
        partners = partners_data['partners']
    
    # Also try to use general images if no specific partner data
    if not partners:
        web_data = getattr(slide, '_web_data', {})
        general_images = web_data.get('images', [])
        if general_images:
            for img in general_images[:3]:
                partners.append({
                    'title': img.get('title', 'Customer'),
                    'image_url': img.get('url') or img.get('download_url')
                })
    
    images_inserted = 0
    for i, placeholder_idx in enumerate(placeholder_indices):
        if i >= len(partners):
            break
        
        partner = partners[i]
        image_url = partner.get('image_url') or partner.get('image') or partner.get('url')
        
        if image_url and image_url.startswith('http'):
            try:
                img_path = download_image_to_temp(image_url, f"customer_{i}.jpg")
                if img_path and os.path.exists(img_path):
                    if fill_image_placeholder(slide, placeholder_idx, img_path):
                        images_inserted += 1
                        continue
            except Exception as e:
                print(f"   ⚠️  Failed to download/insert customer image {i+1}: {e}")
        
        # If image insertion failed, leave placeholder empty or add text
        print(f"   ⚠️  No image available for customer placeholder {placeholder_idx}")
    
    if images_inserted > 0:
        print(f"   ✅ Inserted {images_inserted} customer/partner images")

# ============================================================
# MAIN GENERATION FUNCTION
# ============================================================

def create_presentation(structured_data, output_path):
    """
    Main function to generate presentation from template.
    Falls back to programmatic generation if template not found.
    """
    print(f"🎨 Generating Presentation for {structured_data.get('company_codename')}...")
    
    # Try template-based approach first
    template_path = os.path.join(os.path.dirname(__file__), "templates", "template.pptx")
    if os.path.exists(template_path):
        print("📄 Using template-based generation...")
        success = create_presentation_from_template(structured_data, output_path, template_path)
        if success:
            return
        else:
            print("⚠️  Template generation failed, falling back to programmatic generation...")
    
    # Fallback: Programmatic generation
    print("📝 Using programmatic generation...")
    create_presentation_programmatic(structured_data, output_path)

def create_presentation_from_template(structured_data, output_path, template_path=None):
    """Generates presentation from template by replacing placeholders."""
    global TEMP_IMAGE_DIR
    
    if template_path is None:
        template_path = TEMPLATE_PATH
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return False
    
    print(f"📄 Loading template from: {template_path}")
    
    try:
        prs = Presentation(template_path)
        print(f"   ✅ Template loaded. Found {len(prs.slides)} slides.")
    except Exception as e:
        print(f"❌ Failed to load template: {e}")
        return False
    
    # ============================================================
    # PREPARE DATA DICTIONARY (Match template markers exactly)
    # ============================================================
    data_dict = {
        # Slide 1: Business Overview
        "Business_Overview": format_list(structured_data.get('business_overview', []), max_items=4),
        "Project_Portfolio": format_list(structured_data.get('product_portfolio', []), max_items=8),
        "Applications": format_list(structured_data.get('applications', []), max_items=8),
        "Certifications": format_list(
            [cert.get('name', '') for cert in structured_data.get('web_data', {}).get('certifications', [])],
            max_items=8
        ),
        
        # Slide 2: Financial Metrics
        "EBITDA": clean_text(structured_data.get('financials', {}).get('ebitda', 'N/A')),
        "ROCE": clean_text(structured_data.get('financials', {}).get('roce', 'N/A')),
        "ROE": clean_text(structured_data.get('financials', {}).get('roe', 'N/A')),
        "DEBT": clean_text(structured_data.get('financials', {}).get('debt', 'N/A')),
        "Assumptions": clean_text(structured_data.get('assumptions', 'N/A')),
        "metrics_point": clean_text(structured_data.get('metrics_point', 'N/A')),
        "Upcoming_Facility": clean_text(structured_data.get('upcoming_facility', 'N/A')),
        "Sales": clean_text(structured_data.get('sales', 'N/A')),
        "Global_Presence": clean_text(structured_data.get('global_presence', 'N/A')),
        
        # Slide 3: Investment Highlights
        "Investment_Highlights": format_list(structured_data.get('investment_highlights', []), max_items=5),
    }
    
    # Add any additional financial metrics
    financials = structured_data.get('financials', {})
    for key, value in financials.items():
        marker = key.upper()
        if marker not in data_dict:
            data_dict[marker] = clean_text(value)
    
    # ============================================================
    # FILL PLACEHOLDERS BY INDEX IN ALL SLIDES
    # ============================================================
    print("🔄 Filling placeholders by index...")
    total_replacements = 0
    
    # Debug: List available placeholders for each slide (like index.py)
    for slide_num, slide in enumerate(prs.slides, 1):
        available_placeholders = []
        placeholder_details = []
        # Try to find placeholders by attempting to access common indices
        for idx in range(0, 20):
            try:
                placeholder = slide.placeholders[idx]
                placeholder_info = detect_placeholder_info(slide, idx)
                if placeholder_info:
                    available_placeholders.append(idx)
                    placeholder_details.append(
                        f"idx={placeholder_info['idx']}, type={placeholder_info['type']}, name='{placeholder_info['name']}'"
                    )
            except (KeyError, AttributeError):
                pass
        if available_placeholders:
            print(f"   📋 Slide {slide_num} has placeholders: {sorted(available_placeholders)}")
            for detail in placeholder_details:
                print(f"      - {detail}")
    
    # Skip Slide 1 (index 0) - it's introductory
    print("   📄 Slide 1: Skipping (introductory slide)")
    
    # Slide 2 (index 1): Business Overview
    if len(prs.slides) >= 2:
        slide2 = prs.slides[1]  # Second slide (index 1)
        print("   📄 Slide 2: Business Overview")
        
        # Based on detected placeholders: [10, 11, 12, 13, 14, 15, 16]
        # Index 10: Business Overview text (BODY)
        if fill_text_placeholder(slide2, 10, data_dict.get("Business_Overview", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 10: Business Overview")
        else:
            print("      ⚠️  Placeholder 10 not found in Slide 2")
        
        # Index 14: Product Portfolio (OBJECT)
        if fill_text_placeholder(slide2, 14, data_dict.get("Project_Portfolio", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 14: Product Portfolio")
        else:
            print("      ⚠️  Placeholder 14 not found in Slide 2")
        
        # Index 15: Applications (OBJECT)
        if fill_text_placeholder(slide2, 15, data_dict.get("Applications", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 15: Applications")
        else:
            print("      ⚠️  Placeholder 15 not found in Slide 2")
        
        # Index 16: Certifications (OBJECT)
        if fill_text_placeholder(slide2, 16, data_dict.get("Certifications", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 16: Certifications")
        else:
            print("      ⚠️  Placeholder 16 not found in Slide 2")
        
        # Index 11, 12, 13 are PICTURE placeholders - handled separately for customer images
    
    # Slide 3 (index 2): Financial Metrics
    if len(prs.slides) >= 3:
        slide3 = prs.slides[2]  # Third slide (index 2)
        print("   📄 Slide 3: Financial Metrics")
        
        # Based on detected placeholders: [10, 11, 12, 13, 14, 15]
        # Index 11: Assumptions (BODY)
        if fill_text_placeholder(slide3, 11, data_dict.get("Assumptions", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 11: Assumptions")
        else:
            print("      ⚠️  Placeholder 11 not found in Slide 3")
        
        # Index 14: Metrics point (BODY)
        if fill_text_placeholder(slide3, 14, data_dict.get("metrics_point", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 14: Metrics Point")
        else:
            print("      ⚠️  Placeholder 14 not found in Slide 3")
        
        # Index 15: Upcoming Facility (BODY)
        if fill_text_placeholder(slide3, 15, data_dict.get("Upcoming_Facility", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 15: Upcoming Facility")
        else:
            print("      ⚠️  Placeholder 15 not found in Slide 3")
        
        # Note: Charts (10, 12, 13) would need chart data - skipping for now
        print("      ⚠️  Chart placeholders (10, 12, 13) require chart data - skipping")
    
    # Slide 4 (index 3): Investment Highlights
    if len(prs.slides) >= 4:
        slide4 = prs.slides[3]  # Fourth slide (index 3)
        print("   📄 Slide 4: Investment Highlights")
        
        # Based on detected placeholders: [10]
        # Index 10: Investment Highlights (BODY)
        if fill_text_placeholder(slide4, 10, data_dict.get("Investment_Highlights", "")):
            total_replacements += 1
            print("      ✅ Filled placeholder 10: Investment Highlights")
        else:
            print("      ⚠️  Placeholder 10 not found in Slide 4")
    
    print(f"   ✅ Total replacements: {total_replacements}")
    
    # ============================================================
    # INSERT IMAGES FROM WEB SEARCH
    # ============================================================
    web_data = structured_data.get('web_data', {})
    
    # Store web_data in slide for image insertion functions
    if len(prs.slides) >= 2:
        prs.slides[1]._web_data = web_data
    
    # Insert key customer/partner images into Slide 2 (Business Overview) placeholders (11, 12, 13)
    if len(prs.slides) >= 2:
        slide2 = prs.slides[1]  # Second slide (Business Overview)
        business_info = web_data.get('business_info', {})
        partners = business_info.get('partners', [])
        
        print(f"🖼️  Inserting customer/partner images into Slide 2 placeholders 11, 12, 13...")
        insert_customer_images_by_placeholder(slide2, partners, placeholder_indices=[11, 12, 13])
    
    # ============================================================
    # SAVE PRESENTATION
    # ============================================================
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        unique_output_path = get_unique_output_path(output_path)
        prs.save(unique_output_path)
        print(f"✅ Presentation saved to: {unique_output_path}")
        
        # Clean up temporary images
        if TEMP_IMAGE_DIR and os.path.exists(TEMP_IMAGE_DIR):
            try:
                shutil.rmtree(TEMP_IMAGE_DIR)
            except:
                pass
        
        return True
    except Exception as e:
        print(f"❌ Failed to save presentation: {e}")
        return False

def create_presentation_programmatic(structured_data, output_path):
    """Fallback: Creates presentation programmatically if template not available."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title & Business Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = KELP_INDIGO
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8), Inches(1.2))
    p = title_box.text_frame.add_paragraph()
    p.text = clean_text(structured_data.get("company_codename", "Project Apex")).upper()
    p.font.name = 'Arial'
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = KELP_WHITE
    
    # Business Overview
    ov_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(6), Inches(3.5))
    tf = ov_box.text_frame
    overview = structured_data.get("business_overview", [])
    items = list(overview.values()) if isinstance(overview, dict) else overview
    if isinstance(items, str): items = [items]

    for item in items[:4]:
        p = tf.add_paragraph()
        p.text = f"• {clean_text(item)}"
        p.font.size = Pt(14)
        p.font.color.rgb = KELP_WHITE
        p.space_after = Pt(14)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0), Inches(7.2), Inches(13.33), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.add_paragraph()
    p.text = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
    p.font.name = 'Arial'
    p.font.size = Pt(9)
    p.font.color.rgb = KELP_GREY_TEXT
    p.alignment = PP_ALIGN.CENTER
        
    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    unique_output_path = get_unique_output_path(output_path)
    prs.save(unique_output_path)
    print(f"✅ Success! Presentation saved to: {unique_output_path}")
